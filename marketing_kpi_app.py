import streamlit as st
import pandas as pd
import numpy as np
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import base64

# Page config
st.set_page_config(
    page_title="Marketing KPI Dashboard Generator",
    page_icon="📊",
    layout="wide"
)

def calculate_status(actual, benchmark, direction):
    """Calculate status based on actual vs benchmark and direction"""
    if pd.isna(actual) or pd.isna(benchmark):
        return "Gray"

    if direction == "HigherIsBetter":
        if actual >= benchmark:
            return "Green"
        elif actual >= benchmark * 0.9:
            return "Yellow"
        else:
            return "Red"
    elif direction == "LowerIsBetter":
        if actual <= benchmark:
            return "Green"
        elif actual <= benchmark * 1.1:
            return "Yellow"
        else:
            return "Red"
    return "Gray"

def get_status_color(status):
    """Return RGB color for status"""
    colors = {
        "Green": RGBColor(34, 139, 34),
        "Yellow": RGBColor(255, 165, 0),
        "Red": RGBColor(220, 20, 60),
        "Gray": RGBColor(128, 128, 128)
    }
    return colors.get(status, RGBColor(128, 128, 128))

def create_powerpoint(df, company_name="Your Company"):
    """Create PowerPoint presentation from dataframe"""
    prs = Presentation()

    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"Marketing KPI Dashboard - {company_name} Benchmarks"
    subtitle.text = "Status view across Campaign Types and KPIs"

    # Style title slide with Klaviyo-inspired colors
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)

    # Slide 2: Data table
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Dashboard Status"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Create table
    rows = len(df) + 1  # +1 for header
    cols = 6  # Campaign Type, KPI Name, Benchmark, Actual, Direction, Status

    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)).table

    # Set column widths
    table.columns[0].width = Inches(1.5)  # Campaign Type
    table.columns[1].width = Inches(2.5)  # KPI Name
    table.columns[2].width = Inches(1.2)  # Benchmark
    table.columns[3].width = Inches(1.2)  # Actual
    table.columns[4].width = Inches(1.5)  # Direction
    table.columns[5].width = Inches(1.1)  # Status

    # Header row
    headers = ["Campaign Type", "KPI Name", "Benchmark", "Actual", "Direction", "Status"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(51, 51, 51)

    # Data rows
    for i, row in df.iterrows():
        table.cell(i+1, 0).text = str(row['Campaign Type'])
        table.cell(i+1, 1).text = str(row['KPI Name'])
        table.cell(i+1, 2).text = str(row['Benchmark'])
        table.cell(i+1, 3).text = str(row['Actual'])
        table.cell(i+1, 4).text = str(row['Direction'])

        status = calculate_status(row['Actual'], row['Benchmark'], row['Direction'])
        status_cell = table.cell(i+1, 5)
        status_cell.text = status
        status_cell.fill.solid()
        status_cell.fill.fore_color.rgb = get_status_color(status)
        status_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        status_cell.text_frame.paragraphs[0].font.bold = True

        # Style other cells
        for j in range(5):
            cell = table.cell(i+1, j)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

    return prs

def download_ppt(prs, filename):
    """Create download link for PowerPoint"""
    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)

    b64 = base64.b64encode(ppt_io.read()).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}" download="{filename}">Download PowerPoint</a>'
    return href

# Main app
def main():
    st.title("📊 Marketing KPI Dashboard Generator")
    st.markdown("Generate professional PowerPoint dashboards from your marketing KPI data")

    # Sidebar for company name
    st.sidebar.header("Settings")
    company_name = st.sidebar.text_input("Company Name", value="Your Company")

    # Main content
    tab1, tab2 = st.tabs(["📁 Upload Excel", "✏️ Manual Entry"])

    with tab1:
        st.header("Upload Excel File")
        uploaded_file = st.file_uploader("Choose an Excel file", type=['xlsx', 'xls'])

        if uploaded_file is not None:
            try:
                df = pd.read_excel(uploaded_file)

                # Validate required columns
                required_cols = ['Campaign Type', 'KPI Name', 'Benchmark', 'Actual', 'Direction']
                if all(col in df.columns for col in required_cols):
                    # Filter out rows with NaN in essential columns
                    df_clean = df[required_cols].dropna(subset=['Campaign Type', 'KPI Name'])

                    st.success(f"✅ File uploaded successfully! Found {len(df_clean)} KPIs")

                    # Show preview
                    st.subheader("Data Preview")
                    st.dataframe(df_clean)

                    # Calculate and show status
                    df_clean['Status'] = df_clean.apply(
                        lambda row: calculate_status(row['Actual'], row['Benchmark'], row['Direction']), 
                        axis=1
                    )

                    st.subheader("Status Summary")
                    status_counts = df_clean['Status'].value_counts()
                    col1, col2, col3, col4 = st.columns(4)

                    with col1:
                        st.metric("🟢 Green", status_counts.get('Green', 0))
                    with col2:
                        st.metric("🟡 Yellow", status_counts.get('Yellow', 0))
                    with col3:
                        st.metric("🔴 Red", status_counts.get('Red', 0))
                    with col4:
                        st.metric("⚫ Gray", status_counts.get('Gray', 0))

                    # Generate PowerPoint
                    if st.button("🎯 Generate PowerPoint", type="primary"):
                        with st.spinner("Creating PowerPoint presentation..."):
                            prs = create_powerpoint(df_clean, company_name)

                        st.success("✅ PowerPoint generated successfully!")

                        # Download link
                        filename = f"Marketing_KPI_Dashboard_{company_name.replace(' ', '_')}.pptx"
                        st.markdown(download_ppt(prs, filename), unsafe_allow_html=True)

                else:
                    st.error(f"❌ Missing required columns. Expected: {required_cols}")
                    st.info("Please ensure your Excel file has the correct column headers.")

            except Exception as e:
                st.error(f"❌ Error reading file: {str(e)}")

    with tab2:
        st.header("Manual Data Entry")

        # Initialize session state for manual data
        if 'manual_data' not in st.session_state:
            st.session_state.manual_data = []

        # Form to add new KPI
        with st.form("add_kpi"):
            st.subheader("Add New KPI")
            col1, col2 = st.columns(2)

            with col1:
                campaign_type = st.selectbox("Campaign Type", 
                    ["Events", "Digital/Search", "Influencer", "Commercial", "Social", "Other"])
                kpi_name = st.text_input("KPI Name")
                benchmark = st.number_input("Benchmark", value=0.0)

            with col2:
                actual = st.number_input("Actual", value=0.0)
                direction = st.selectbox("Direction", ["HigherIsBetter", "LowerIsBetter"])

            if st.form_submit_button("➕ Add KPI"):
                if kpi_name:
                    new_kpi = {
                        'Campaign Type': campaign_type,
                        'KPI Name': kpi_name,
                        'Benchmark': benchmark,
                        'Actual': actual,
                        'Direction': direction
                    }
                    st.session_state.manual_data.append(new_kpi)
                    st.success(f"✅ Added: {kpi_name}")
                else:
                    st.error("❌ Please enter a KPI name")

        # Display current data
        if st.session_state.manual_data:
            st.subheader("Current KPIs")
            df_manual = pd.DataFrame(st.session_state.manual_data)

            # Calculate status
            df_manual['Status'] = df_manual.apply(
                lambda row: calculate_status(row['Actual'], row['Benchmark'], row['Direction']), 
                axis=1
            )

            st.dataframe(df_manual)

            col1, col2 = st.columns(2)
            with col1:
                if st.button("🗑️ Clear All Data"):
                    st.session_state.manual_data = []
                    st.rerun()

            with col2:
                if st.button("🎯 Generate PowerPoint", type="primary", key="manual_ppt"):
                    with st.spinner("Creating PowerPoint presentation..."):
                        prs = create_powerpoint(df_manual, company_name)

                    st.success("✅ PowerPoint generated successfully!")

                    # Download link
                    filename = f"Marketing_KPI_Dashboard_{company_name.replace(' ', '_')}.pptx"
                    st.markdown(download_ppt(prs, filename), unsafe_allow_html=True)
        else:
            st.info("👆 Add some KPIs using the form above to get started!")

if __name__ == "__main__":
    main()
